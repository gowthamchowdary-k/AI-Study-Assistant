import os
from pypdf import PdfReader
import docx
from pptx import Presentation
from google import genai
from google.genai import types
from config import GEMINI_API_KEY, GEMINI_MODEL

def read_pdf(file_path):
    """
    Reads a PDF page by page and returns text and page numbers.
    """
    reader = PdfReader(file_path)
    pages = []
    file_name = os.path.basename(file_path)

    for page_number, page in enumerate(reader.pages, start=1):
        extracted = page.extract_text()
        if extracted and extracted.strip():
            pages.append({
                "file": file_name,
                "page": page_number,
                "text": extracted
            })
    return pages

def read_docx(file_path):
    """
    Reads a Word Document (.docx) and returns it as a single page.
    """
    doc = docx.Document(file_path)
    file_name = os.path.basename(file_path)
    
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)
            
    text = "\n".join(full_text)
    if not text.strip():
        return []
        
    return [{
        "file": file_name,
        "page": 1,
        "text": text
    }]

def read_pptx(file_path):
    """
    Reads a PowerPoint Presentation (.pptx) slide by slide.
    """
    prs = Presentation(file_path)
    pages = []
    file_name = os.path.basename(file_path)

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        text = "\n".join(slide_text)
        if text.strip():
            pages.append({
                "file": file_name,
                "page": slide_number,
                "text": text
            })
    return pages

def read_txt(file_path):
    """
    Reads a plain text file (.txt).
    """
    file_name = os.path.basename(file_path)
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as f:
            text = f.read()

    if not text.strip():
        return []
        
    return [{
        "file": file_name,
        "page": 1,
        "text": text
    }]

def read_image_ocr(file_path):
    """
    Uses Gemini API to perform OCR on images.
    Supports PNG, JPG, JPEG, WEBP.
    """
    if not GEMINI_API_KEY:
        raise ValueError("GEMINI_API_KEY is not set in your .env file. Image OCR requires a Gemini API key.")

    file_name = os.path.basename(file_path)
    ext = os.path.splitext(file_path)[1].lower()
    
    mime_type = "image/png"
    if ext in [".jpg", ".jpeg"]:
        mime_type = "image/jpeg"
    elif ext == ".webp":
        mime_type = "image/webp"

    with open(file_path, "rb") as f:
        image_bytes = f.read()

    try:
        client = genai.Client(api_key=GEMINI_API_KEY)
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=[
                types.Part.from_bytes(data=image_bytes, mime_type=mime_type),
                "Perform OCR on this image. Extract all text verbatim and preserve paragraphs. Do not add any summary, explanation, or comments. Just return the text found."
            ]
        )
        text = response.text or ""
        if not text.strip():
            return []
            
        return [{
            "file": file_name,
            "page": 1,
            "text": text
        }]
    except Exception as e:
        print(f"IMAGE OCR ERROR: Failed to parse image using Gemini. Details: {e}")
        raise RuntimeError(f"OCR failed for image {file_name}: {e}")

def read_document(file_path):
    """
    Unified entrypoint to read any supported document.
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return read_pdf(file_path)
    elif ext == ".docx":
        return read_docx(file_path)
    elif ext == ".pptx":
        return read_pptx(file_path)
    elif ext == ".txt":
        return read_txt(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
        return read_image_ocr(file_path)
    else:
        raise ValueError(f"Unsupported file extension '{ext}'")